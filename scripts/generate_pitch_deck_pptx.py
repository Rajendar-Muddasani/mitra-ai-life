from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


OUTPUT_PATH = Path(__file__).resolve().parents[1] / "site" / "mitra-ai-life-education-pitch-deck.pptx"
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

BG = RGBColor(4, 7, 18)
BG_ALT = RGBColor(9, 17, 37)
PANEL = RGBColor(14, 24, 56)
PANEL_STRONG = RGBColor(24, 42, 94)
TEXT = RGBColor(245, 248, 255)
MUTED = RGBColor(159, 179, 218)
BLUE = RGBColor(83, 166, 255)
CYAN = RGBColor(104, 225, 255)
INDIGO = RGBColor(121, 145, 255)


def add_background(slide, accent=False):
    background = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        0,
        0,
        SLIDE_WIDTH,
        SLIDE_HEIGHT,
    )
    background.fill.solid()
    background.fill.fore_color.rgb = BG_ALT if accent else BG
    background.line.fill.background()

    glow_left = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.OVAL,
        Inches(-0.4),
        Inches(-0.25),
        Inches(3.0),
        Inches(3.0),
    )
    glow_left.fill.solid()
    glow_left.fill.fore_color.rgb = BLUE
    glow_left.fill.transparency = 0.78
    glow_left.line.fill.background()

    glow_right = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.OVAL,
        Inches(10.7),
        Inches(-0.35),
        Inches(2.6),
        Inches(2.6),
    )
    glow_right.fill.solid()
    glow_right.fill.fore_color.rgb = CYAN if accent else INDIGO
    glow_right.fill.transparency = 0.84
    glow_right.line.fill.background()


def add_textbox(slide, left, top, width, height, text, *, font_size=18, color=TEXT, bold=False,
                font_name="Aptos", align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.vertical_anchor = MSO_ANCHOR.TOP
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    run = paragraph.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_title_block(slide, eyebrow, title, subtitle=None):
    add_textbox(slide, Inches(0.75), Inches(0.5), Inches(3.4), Inches(0.35), eyebrow.upper(), font_size=11, color=CYAN)
    add_textbox(slide, Inches(0.75), Inches(0.9), Inches(8.0), Inches(1.5), title, font_size=26, bold=True, font_name="Aptos Display")
    if subtitle:
      add_textbox(slide, Inches(9.15), Inches(0.95), Inches(3.1), Inches(1.0), subtitle, font_size=12, color=MUTED)


def add_card(slide, left, top, width, height, title, body, *, strong=False, glow=False, bullet_list=False):
    card = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = PANEL_STRONG if strong else PANEL
    if glow:
        card.line.color.rgb = CYAN
        card.line.transparency = 0.35
    else:
        card.line.color.rgb = BLUE
        card.line.transparency = 0.6

    title_box = add_textbox(slide, left + Inches(0.18), top + Inches(0.12), width - Inches(0.36), Inches(0.4), title,
                            font_size=16, bold=True, font_name="Aptos Display")
    title_box.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT

    body_box = slide.shapes.add_textbox(left + Inches(0.18), top + Inches(0.55), width - Inches(0.36), height - Inches(0.68))
    frame = body_box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.vertical_anchor = MSO_ANCHOR.TOP

    if bullet_list:
        for index, item in enumerate(body):
            paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
            paragraph.text = item
            paragraph.level = 0
            paragraph.bullet = True
            paragraph.font.name = "Aptos"
            paragraph.font.size = Pt(11)
            paragraph.font.color.rgb = MUTED
            paragraph.space_after = Pt(6)
    else:
        paragraph = frame.paragraphs[0]
        paragraph.text = body
        paragraph.font.name = "Aptos"
        paragraph.font.size = Pt(11.5)
        paragraph.font.color.rgb = MUTED


def add_pill_row(slide, pills):
    left = Inches(0.75)
    top = Inches(4.1)
    for pill in pills:
        width = max(1.15, min(2.15, 0.48 + len(pill) * 0.07))
        shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, Inches(width), Inches(0.42))
        shape.fill.solid()
        shape.fill.fore_color.rgb = PANEL
        shape.line.color.rgb = CYAN
        shape.line.transparency = 0.5
        add_textbox(slide, left + Inches(0.08), top + Inches(0.06), Inches(width) - Inches(0.16), Inches(0.2), pill,
                    font_size=10.5, color=TEXT)
        left += Inches(width + 0.12)


def build_cover(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, accent=True)
    add_textbox(slide, Inches(0.75), Inches(0.5), Inches(2.5), Inches(0.3), "FOUNDER DECK", font_size=11, color=CYAN)
    add_textbox(slide, Inches(0.75), Inches(0.95), Inches(6.7), Inches(1.8), "Mitra AI Life Education",
                font_size=30, bold=True, font_name="Aptos Display")
    add_textbox(slide, Inches(0.75), Inches(2.05), Inches(6.7), Inches(1.1),
                "Practical AI learning for daily life in India, built English first and translated into Telugu next.",
                font_size=14, color=MUTED)
    add_pill_row(slide, ["Beginner first", "English source", "Telugu next", "Free-first launch", "Visual lessons"])
    add_card(slide, Inches(8.15), Inches(1.0), Inches(4.3), Inches(1.35), "Level 1 becomes the first hour",
             "Mostly watch, listen, and feel the value before asking users to do much.", strong=True, glow=True)
    add_card(slide, Inches(8.15), Inches(2.55), Inches(4.3), Inches(1.35), "Level 2 builds confidence",
             "Hands-on daily-life wins turn curiosity into habit and prepare users for paid depth.")
    add_card(slide, Inches(8.15), Inches(4.1), Inches(4.3), Inches(1.35), "Premium lane stays separate",
             "Engineering project services remain distinct from mass-market beginner education.")


def build_brand(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_title_block(slide, "Brand System", "One clean naming stack across product, repo, and future company setup.")
    add_card(slide, Inches(0.75), Inches(1.8), Inches(5.8), Inches(3.1), "Public-facing brand",
             "Mitra AI Life Education\n\nWarm, clear, Indian-friendly, and flexible enough for lessons, community, and future services.",
             strong=True)
    add_card(slide, Inches(6.8), Inches(1.8), Inches(5.75), Inches(3.1), "Working system names",
             [
                 "Workspace: Mitra_AI_Life",
                 "Repository: mitra-ai-life",
                 "Website title: Mitra AI Life Education",
                 "Future company: Mitra AI Life Education Private Limited",
             ],
             bullet_list=True,
             glow=True)


def build_why(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, accent=True)
    add_title_block(slide, "Why This Can Win", "Most AI learning still starts too technical and too cold for the real mass market.",
                    "The opportunity is becoming the trusted guide for parents, students, homemakers, and first-time workers.")
    add_card(slide, Inches(0.75), Inches(2.1), Inches(3.85), Inches(2.3), "Gap",
             "Most AI education is jargon-heavy, desktop-first, and disconnected from Indian daily routines.", glow=True)
    add_card(slide, Inches(4.75), Inches(2.1), Inches(3.85), Inches(2.3), "Need",
             "People want help with messages, learning, planning, budgeting, travel, and confidence, not model theory.")
    add_card(slide, Inches(8.75), Inches(2.1), Inches(3.8), Inches(2.3), "Advantage",
             "English-first content with Telugu adaptation creates a practical lane that is easier to trust and scale.")


def build_first_hour(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_title_block(slide, "First-Hour Strategy", "Level 1 should feel like a guided show, not a rushed tutorial.")
    add_card(slide, Inches(0.75), Inches(1.95), Inches(3.9), Inches(2.0), "01  Wonder",
             "Open with polished stories and voice-led demos so the learner feels the possibility of AI in normal life.",
             strong=True)
    add_card(slide, Inches(4.75), Inches(1.95), Inches(3.9), Inches(2.0), "02  Safety",
             "Explain what AI does well, where it can be wrong, and why human checking matters from day one.")
    add_card(slide, Inches(8.75), Inches(1.95), Inches(3.8), Inches(2.0), "03  Tiny action",
             "Ask for only one or two simple actions near the end, then point clearly into the next level.")
    add_card(slide, Inches(0.75), Inches(4.25), Inches(5.75), Inches(2.1), "Level 1: First Step",
             "45 to 60 minutes\n\nMostly view, listen, and observe. Strong story flow. Minimal typing. Goal: trust and curiosity.",
             glow=True)
    add_card(slide, Inches(6.8), Inches(4.25), Inches(5.75), Inches(2.1), "Level 2: Daily Help",
             "60 to 90 minutes\n\nReal hands-on tasks across messaging, study help, planning, and bilingual drafting. Goal: confidence and repeat use.",
             strong=True)


def build_ladder(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, accent=True)
    add_title_block(slide, "Curriculum Ladder", "Free trust first, paid depth later.")
    cards = [
        ("First Step", "45 to 60 min", "Free"),
        ("Daily Help", "60 to 90 min", "Free"),
        ("Smart Basics", "1 hour", "100 INR"),
        ("Work Smart", "2 hours", "200 INR"),
        ("Life Upgrade", "3 hours", "300 INR"),
        ("Power User", "4 hours", "400 INR"),
        ("Build With AI", "5 hours", "500 INR"),
    ]
    positions = [
        (0.75, 2.1),
        (2.55, 2.1),
        (4.35, 2.1),
        (6.15, 2.1),
        (7.95, 2.1),
        (9.75, 2.1),
        (11.55, 2.1),
    ]
    for index, ((title, duration, price), (left, top)) in enumerate(zip(cards, positions)):
        card = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(1.55), Inches(2.5))
        card.fill.solid()
        card.fill.fore_color.rgb = PANEL_STRONG if index < 2 else PANEL
        card.line.color.rgb = CYAN if index == 6 else BLUE
        card.line.transparency = 0.45
        add_textbox(slide, Inches(left + 0.12), Inches(top + 0.18), Inches(1.3), Inches(0.52), title, font_size=11.5,
                    bold=True, font_name="Aptos Display")
        add_textbox(slide, Inches(left + 0.12), Inches(top + 1.02), Inches(1.3), Inches(0.38), duration, font_size=10.5, color=MUTED)
        add_textbox(slide, Inches(left + 0.12), Inches(top + 1.72), Inches(1.3), Inches(0.42), price, font_size=11.5, color=TEXT, bold=True)
    add_card(slide, Inches(0.75), Inches(5.1), Inches(11.8), Inches(0.9), "Starter promise",
             "Level 1 builds trust, Level 2 builds confidence, and Level 3 becomes the first paid upgrade.", glow=True)


def build_experience(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_title_block(slide, "Lesson Experience", "Teach one topic in multiple easy-to-consume formats.")
    add_card(slide, Inches(0.75), Inches(2.0), Inches(3.85), Inches(3.25), "Per lesson pack",
             [
                 "anchor story",
                 "comic or cartoon explainer",
                 "voice-led slide sequence",
                 "cheat sheet and recap page",
                 "worksheet and challenge",
             ],
             bullet_list=True)
    add_card(slide, Inches(4.75), Inches(2.0), Inches(3.85), Inches(3.25), "Production workflow",
             [
                 "write the master in simple English",
                 "generate visuals and voice assets",
                 "adapt meaning-first into Telugu",
                 "run human review before publishing",
             ],
             bullet_list=True,
             strong=True,
             glow=True)
    add_card(slide, Inches(8.75), Inches(2.0), Inches(3.8), Inches(3.25), "Trust guardrails",
             [
                 "no exaggerated claims",
                 "no blind trust in AI output",
                 "clear legal and safety disclaimers",
                 "separate premium services clearly",
             ],
             bullet_list=True)


def build_gtm(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, accent=True)
    add_title_block(slide, "Go-To-Market", "Publish small, measure honestly, then scale.")
    phases = [
        ("Phase 0", "Finalize brand, repo, landing message, and first free lesson structure."),
        ("Phase 1", "Launch Levels 1 and 2 free, collect feedback, and validate beginner understanding."),
        ("Phase 2", "Introduce paid Levels 3 to 5, combo pricing, progress tracking, and certificates."),
        ("Phase 3+", "Expand segments, then keep engineering project studio as a separate premium lane."),
    ]
    x_positions = [0.75, 3.75, 6.75, 9.75]
    for index, ((title, body), left) in enumerate(zip(phases, x_positions)):
        add_card(slide, Inches(left), Inches(2.2), Inches(2.6), Inches(3.0), title, body, strong=index == 0, glow=index == 2)


def build_business(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_title_block(slide, "Business Shape", "Free acquisition feeds paid practical learning and premium services later.")
    add_card(slide, Inches(0.75), Inches(2.05), Inches(3.85), Inches(2.45), "Free engine",
             "Level 1 and Level 2 drive trust, reach, and word of mouth before payments are introduced.", glow=True)
    add_card(slide, Inches(4.75), Inches(2.05), Inches(3.85), Inches(2.45), "Paid curriculum",
             "Levels 3 to 7 sell deeper results: communication, productivity, structured prompting, and creation.")
    add_card(slide, Inches(8.75), Inches(2.05), Inches(3.8), Inches(2.45), "Premium studio",
             "Engineering projects, implementation support, and institution-backed builds operate as a distinct high-support business.")
    add_card(slide, Inches(0.75), Inches(4.95), Inches(3.85), Inches(1.05), "Launch objective", "100 early users", strong=True)
    add_card(slide, Inches(4.75), Inches(4.95), Inches(3.85), Inches(1.05), "Key free metric", "Level 1 completion")
    add_card(slide, Inches(8.75), Inches(4.95), Inches(3.8), Inches(1.05), "Key paid metric", "Free-to-paid conversion", glow=True)


def build_platform(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, accent=True)
    add_title_block(slide, "Platform Shape", "Keep source content in GitHub. Keep heavy media in object storage.",
                    "The repo should own lesson text, prompts, manifests, scripts, and site code.")
    add_card(slide, Inches(0.75), Inches(2.15), Inches(5.75), Inches(3.45), "Source layer",
             [
                 "Markdown lesson source",
                 "Prompt library and manifests",
                 "Planning docs and scripts",
                 "Static website code",
             ],
             bullet_list=True,
             strong=True,
             glow=True)
    add_card(slide, Inches(6.8), Inches(2.15), Inches(5.75), Inches(3.45), "Delivery layer",
             [
                 "Cloudflare Pages or Vercel",
                 "Cloudflare R2 or S3-compatible storage",
                 "PostHog analytics",
                 "Supabase for auth and progress later",
             ],
             bullet_list=True)


def build_close(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_title_block(slide, "Immediate Build", "Ship the first impression before building the full platform.")
    add_card(slide, Inches(0.75), Inches(2.3), Inches(3.85), Inches(2.8), "Now",
             "Use the finalized brand stack, launch the black-and-blue founder deck, and shape Level 1 as a cinematic first hour.",
             strong=True)
    add_card(slide, Inches(4.75), Inches(2.3), Inches(3.85), Inches(2.8), "Next",
             "Create the free lesson assets in English, then move into Telugu adaptation and feedback capture.", glow=True)
    add_card(slide, Inches(8.75), Inches(2.3), Inches(3.8), Inches(2.8), "Later",
             "Add payments, user progress, and premium project offerings only after the free experience proves demand.")


def main():
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    build_cover(prs)
    build_brand(prs)
    build_why(prs)
    build_first_hour(prs)
    build_ladder(prs)
    build_experience(prs)
    build_gtm(prs)
    build_business(prs)
    build_platform(prs)
    build_close(prs)

    OUTPUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUTPUT_PATH)
    print(f"Wrote {OUTPUT_PATH}")


if __name__ == "__main__":
    main()