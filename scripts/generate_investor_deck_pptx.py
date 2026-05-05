"""
Investor pitch deck for Mitra AI Life Education.
10 slides — market-size focused, data-driven, clean.
Distinct from the founder deck (pitch-deck.html / generate_pitch_deck_pptx.py).
"""
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


OUTPUT_PATH = Path(__file__).resolve().parents[1] / "site" / "mitra-ai-life-education-investor-deck.pptx"
SLIDE_WIDTH  = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

# ── colour system (same brand palette) ──────────────────────────────────────
BG           = RGBColor(4,   7,   18)
BG_ALT       = RGBColor(9,   17,  37)
PANEL        = RGBColor(14,  24,  56)
PANEL_STRONG = RGBColor(24,  42,  94)
TEXT         = RGBColor(245, 248, 255)
MUTED        = RGBColor(159, 179, 218)
BLUE         = RGBColor(83,  166, 255)
CYAN         = RGBColor(104, 225, 255)
INDIGO       = RGBColor(121, 145, 255)
GREEN        = RGBColor(72,  219, 150)
AMBER        = RGBColor(255, 196, 82)


# ── primitives ───────────────────────────────────────────────────────────────

def add_bg(slide, accent=False):
    bg = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT)
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG_ALT if accent else BG
    bg.line.fill.background()

    def glow(l, t, w, h, color, opacity):
        s = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.OVAL, Inches(l), Inches(t), Inches(w), Inches(h))
        s.fill.solid()
        s.fill.fore_color.rgb = color
        s.fill.transparency = opacity
        s.line.fill.background()

    glow(-0.4, -0.25, 3.0, 3.0, BLUE, 0.80)
    glow(10.7, -0.35, 2.6, 2.6, CYAN if accent else INDIGO, 0.86)


def tb(slide, l, t, w, h, text, *, pt=14, color=TEXT, bold=False,
       font="Aptos", align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.vertical_anchor = MSO_ANCHOR.TOP
    p = frame.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(pt)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def title_block(slide, eyebrow, title, sub=None):
    tb(slide, 0.75, 0.5,  3.4, 0.35, eyebrow.upper(), pt=10, color=CYAN)
    tb(slide, 0.75, 0.85, 8.0, 1.5,  title,           pt=25, bold=True, font="Aptos Display")
    if sub:
        tb(slide, 9.15, 0.9, 3.6, 1.0, sub, pt=11, color=MUTED)


def card(slide, l, t, w, h, headline, body, *,
         strong=False, accent_color=None, bullet=False, pt=11):
    c = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(l), Inches(t), Inches(w), Inches(h))
    c.fill.solid()
    c.fill.fore_color.rgb = PANEL_STRONG if strong else PANEL
    line_color = accent_color if accent_color else (CYAN if strong else BLUE)
    c.line.color.rgb = line_color
    c.line.transparency = 0.35 if strong else 0.60

    tb(slide, l + 0.18, t + 0.12, w - 0.36, 0.42,
       headline, pt=14, bold=True, font="Aptos Display")

    body_box = slide.shapes.add_textbox(
        Inches(l + 0.18), Inches(t + 0.55), Inches(w - 0.36), Inches(h - 0.68))
    frame = body_box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.vertical_anchor = MSO_ANCHOR.TOP

    if bullet:
        for i, item in enumerate(body):
            para = frame.paragraphs[0] if i == 0 else frame.add_paragraph()
            para.text = item
            para.bullet = True
            para.font.name = "Aptos"
            para.font.size = Pt(pt)
            para.font.color.rgb = MUTED
            para.space_after = Pt(5)
    else:
        para = frame.paragraphs[0]
        para.text = body
        para.font.name = "Aptos"
        para.font.size = Pt(pt)
        para.font.color.rgb = MUTED


def metric(slide, l, t, w, h, value, label, color=CYAN):
    """Large number + small label block."""
    card(slide, l, t, w, h, "", "", strong=False)
    tb(slide, l + 0.15, t + 0.18, w - 0.3, 0.9, value,
       pt=26, bold=True, color=color, font="Aptos Display", align=PP_ALIGN.CENTER)
    tb(slide, l + 0.15, t + 1.05, w - 0.3, 0.55, label,
       pt=10, color=MUTED, align=PP_ALIGN.CENTER)


def divider(slide, y):
    line = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.75), Inches(y), Inches(11.85), Pt(1))
    line.fill.solid()
    line.fill.fore_color.rgb = BLUE
    line.fill.transparency = 0.70
    line.line.fill.background()


# ── slides ───────────────────────────────────────────────────────────────────

def build_cover(prs):
    """Slide 1 — Cover: market hook + one-line thesis."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, accent=True)

    tb(slide, 0.75, 0.5,  3.0, 0.3, "INVESTOR BRIEF", pt=10, color=CYAN)
    tb(slide, 0.75, 0.9,  7.2, 1.8,
       "Mitra AI Life Education",
       pt=30, bold=True, font="Aptos Display")
    tb(slide, 0.75, 2.1,  7.2, 0.6,
       "Seed round  |  India AI literacy, English first, Telugu next",
       pt=13, color=INDIGO, bold=True)
    tb(slide, 0.75, 2.85, 7.0, 0.9,
       "300 million internet-connected Indians have never used an AI tool. "
       "We are building the simplest, most relatable first-hour AI experience "
       "for them — in their language, with their everyday problems.",
       pt=12, color=MUTED)

    metric(slide, 0.75, 4.0,  2.7, 1.85, "₹6B+",    "India edtech market 2024",  CYAN)
    metric(slide, 3.65, 4.0,  2.7, 1.85, "300M+",   "AI-untouched internet users", BLUE)
    metric(slide, 6.55, 4.0,  2.7, 1.85, "30%",     "Online education CAGR 2024–30", INDIGO)
    metric(slide, 9.45, 4.0,  2.7, 1.85, "₹1 Cr",  "Seed ask", GREEN)


def build_problem(prs):
    """Slide 2 — Problem: the AI literacy gap in India."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    title_block(slide, "The Problem",
                "300 million connected Indians have no trusted first step into AI.",
                sub="70% of Indian internet users have never used an AI tool.")

    card(slide, 0.75, 2.05, 3.85, 2.85,
         "The content gap",
         "Existing AI education is English-heavy, desktop-first, jargon-filled, and built for tech-savvy users. "
         "It does not speak to homemakers, small shopkeepers, or first-generation smartphone owners.",
         )

    card(slide, 4.75, 2.05, 3.85, 2.85,
         "The language gap",
         "~600 million Indians use a regional language as their primary digital language. "
         "Telugu alone covers 90M+ people. Almost zero structured AI learning exists in Telugu.",
         strong=True)

    card(slide, 8.75, 2.05, 3.8, 2.85,
         "The trust gap",
         "When people do encounter AI, it feels like magic they cannot control — or a risk they do not understand. "
         "There is no guide that is warm, simple, and honest.",
         accent_color=AMBER)

    divider(slide, 5.3)
    tb(slide, 0.75, 5.45, 11.85, 0.5,
       "Nobody is building the reliable, beginner-first, vernacular AI learning path for the Indian mass market.",
       pt=12, color=TEXT, bold=True, align=PP_ALIGN.CENTER)


def build_solution(prs):
    """Slide 3 — Solution: what we are building."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, accent=True)
    title_block(slide, "Our Solution",
                "A beginner-first AI learning platform built for daily Indian life.",
                sub="Practical outcomes first. Telugu as the first translation.")

    card(slide, 0.75, 2.05, 5.75, 3.5,
         "What the learner gets",
         [
             "Relatable stories drawn from Indian daily life (not Silicon Valley demos)",
             "Watch-first, act-later lesson structure — low anxiety, high trust",
             "Simple English source, then meaning-first Telugu translation",
             "Comic strips, voice-led sequences, cheat sheets, and micro-worksheets",
             "Human-reviewed content — never publish raw AI output",
         ],
         bullet=True, strong=True)

    card(slide, 6.8, 2.05, 5.75, 3.5,
         "What makes us different",
         [
             "Vernacular-first roadmap from day one, not an afterthought",
             "Free first two levels — trust before transactions",
             "Production workflow: English master → auto-draft → Telugu human review",
             "Content platform is separate from premium engineering project studio",
             "Designed for WhatsApp-era learners with a smartphone, not a PC",
         ],
         bullet=True)

    divider(slide, 5.95)
    tb(slide, 0.75, 6.1, 11.85, 0.5,
       "We are the trusted Mitra (friend) for the first million Indians to use AI confidently.",
       pt=12, color=CYAN, bold=True, align=PP_ALIGN.CENTER)


def build_market(prs):
    """Slide 4 — Market size: TAM / SAM / SOM."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    title_block(slide, "Market Size",
                "India edtech is a ₹50,000 Cr opportunity. We are targeting a fast-growing vernacular AI slice.",
                sub="India online education 2024 → 2030 CAGR: 30%")

    # TAM / SAM / SOM funnel cards
    card(slide, 0.75, 2.0,  3.85, 3.6,
         "TAM — India online education",
         "₹50,000 Cr+ (2024)\nAll Indians learning online: students, working adults, homemakers, seniors.\n\n"
         "AI literacy will become a required skill across every segment.",
         strong=True)

    card(slide, 4.75, 2.0,  3.85, 3.6,
         "SAM — Regional AI literacy",
         "₹6,500 Cr (estimated)\nInternet-connected, semi-urban and rural Indians who will need AI skill "
         "training in their own language within the next 5 years.\n\n"
         "Telugu, Hindi, Tamil, Kannada, Marathi lead.",
         accent_color=BLUE)

    card(slide, 8.75, 2.0,  3.8,  3.6,
         "SOM — Telugu + Hindi first (3 yr)",
         "₹350 Cr reachable\n500K paid learners × ₹700 average annual basket.\n\n"
         "Capturing 0.2% of the SAM in 3 years via organic WhatsApp distribution, community referrals, "
         "and school partnerships.",
         accent_color=GREEN)

    divider(slide, 5.95)
    tb(slide, 0.75, 6.1, 11.85, 0.55,
       "India is adding 30M new internet users per year — most of them mobile-first, language-first, AI-ready by 2026.",
       pt=11.5, color=MUTED, align=PP_ALIGN.CENTER)


def build_product(prs):
    """Slide 5 — Product: level ladder and roadmap."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, accent=True)
    title_block(slide, "Product Roadmap",
                "Seven levels. Free trust first. Paid depth from Level 3.")

    # 7 level mini-cards
    levels = [
        ("L1", "First Step",   "45–60 min", "FREE",   CYAN,   True),
        ("L2", "Daily Help",   "60–90 min", "FREE",   CYAN,   True),
        ("L3", "Smart Basics", "1 hr",      "₹100",   BLUE,   False),
        ("L4", "Work Smart",   "2 hrs",     "₹200",   BLUE,   False),
        ("L5", "Life Upgrade", "3 hrs",     "₹300",   INDIGO, False),
        ("L6", "Power User",   "4 hrs",     "₹400",   INDIGO, False),
        ("L7", "Build With AI","5 hrs",     "₹500",   GREEN,  False),
    ]
    card_w = 1.55
    gap    = 0.1
    start  = 0.75
    for i, (code, name, dur, price, col, is_free) in enumerate(levels):
        lx = start + i * (card_w + gap)
        c = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            Inches(lx), Inches(2.05), Inches(card_w), Inches(2.85))
        c.fill.solid()
        c.fill.fore_color.rgb = PANEL_STRONG if is_free else PANEL
        c.line.color.rgb = col
        c.line.transparency = 0.35 if is_free else 0.55
        tb(slide, lx + 0.1, 2.18, card_w - 0.2, 0.3,  code,  pt=9,  color=col,  bold=True)
        tb(slide, lx + 0.1, 2.48, card_w - 0.2, 0.48, name,  pt=11, color=TEXT, bold=True, font="Aptos Display")
        tb(slide, lx + 0.1, 2.98, card_w - 0.2, 0.36, dur,   pt=9.5, color=MUTED)
        tb(slide, lx + 0.1, 3.52, card_w - 0.2, 0.42, price, pt=12, color=col, bold=True)

    card(slide, 0.75, 5.25, 5.75, 1.3,
         "Combo bundles available from Phase 2",
         "L3+L4+L5 bundle at ₹450 and full ladder bundle at ₹1,200 drive higher average basket.",
         strong=True)
    card(slide, 6.80, 5.25, 5.75, 1.3,
         "Premium engineering studio is separate",
         "High-support implementation tracks for teams, colleges, and institutions — different pricing, different support SLA.")


def build_traction(prs):
    """Slide 6 — Traction plan / milestones (pre-revenue, honest)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    title_block(slide, "Go-To-Market & Milestones",
                "Four phases. Measure completion before conversion.",
                sub="Distribution: WhatsApp, YouTube, vernacular community groups, school tie-ups.")

    phases = [
        ("Phase 0\nNow – M3",
         "Brand locked, GitHub launched, pitch deck live.\n"
         "Level 1 script, storyboard, and interactive HTML built.\n"
         "Invite 50 beta testers from Hyderabad network.",
         CYAN, True),
        ("Phase 1\nM4 – M6",
         "Level 1 and Level 2 fully live (free).\n"
         "1,000 learners, 60%+ Level 1 completion rate.\n"
         "Telugu version published and reviewed.",
         BLUE, False),
        ("Phase 2\nM7 – M12",
         "Levels 3–5 go live with pricing.\n"
         "10,000 learners; 10% paid conversion = 1,000 paying users.\n"
         "₹50L ARR. Combo bundles, certificates launched.",
         INDIGO, False),
        ("Phase 3\nM13 – M18",
         "50,000+ learners. ₹1.5 Cr ARR.\n"
         "Second language (Hindi or Tamil) adapted.\n"
         "Series A readiness: 3x revenue growth proof.",
         GREEN, False),
    ]
    for i, (title, body, col, strong_flag) in enumerate(phases):
        lx = 0.75 + i * 3.15
        card(slide, lx, 2.1, 2.85, 3.45, title, body,
             strong=strong_flag, accent_color=col)

    divider(slide, 5.9)
    tb(slide, 0.75, 6.05, 11.85, 0.55,
       "No paid content until Level 1 completion rate exceeds 60% — quality gates before monetisation.",
       pt=11.5, color=AMBER, bold=True, align=PP_ALIGN.CENTER)


def build_unit_economics(prs):
    """Slide 7 — Business model + unit economics."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, accent=True)
    title_block(slide, "Business Model & Unit Economics",
                "Low CAC. High LTV. Free acquisition → paid depth.")

    card(slide, 0.75, 2.05, 3.85, 3.55,
         "Revenue streams",
         [
             "Level 3–7 individual purchases (₹100–₹500)",
             "Combo bundles (₹450 and ₹1,200 all-access)",
             "Annual subscription — planned Year 2",
             "Institution/school batch licensing",
             "Premium engineering project studio (separate P&L)",
         ],
         bullet=True, strong=True)

    card(slide, 4.75, 2.05, 3.85, 3.55,
         "Unit economics (target Year 2)",
         [
             "CAC target: ₹40 (WhatsApp organic)",
             "Avg first-year basket: ₹800",
             "LTV (3-year): ₹2,400",
             "LTV / CAC: 60x",
             "Gross margin target: 72%",
             "Level 1 completion → paid conversion: 12–18%",
         ],
         bullet=True, accent_color=GREEN)

    card(slide, 8.75, 2.05, 3.8, 3.55,
         "Why margins stay high",
         [
             "Content created once, served millions",
             "Markdown source → multi-format auto-export",
             "Object storage (Cloudflare R2) < $10/month at 100K users",
             "No live instructors in core curriculum",
             "Human review only at final publication step",
         ],
         bullet=True)


def build_revenue(prs):
    """Slide 8 — 3-year revenue projections."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    title_block(slide, "Revenue Projections",
                "Conservative base case. No institutional revenue included.",
                sub="All figures in Indian Rupees. ARR = annualised recurring revenue.")

    # Projection metrics by year
    years = [
        ("Year 1",  "50K",  "5,000",   "₹40L",   "₹0.40 Cr"),
        ("Year 2",  "200K", "30,000",  "₹2.4 Cr","₹2.4 Cr"),
        ("Year 3",  "600K", "100,000", "₹8 Cr",  "₹8 Cr"),
    ]
    col_colors = [BLUE, INDIGO, GREEN]
    for i, (yr, learners, paid, arrint, arr) in enumerate(years):
        lx = 0.75 + i * 4.1
        c = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            Inches(lx), Inches(2.0), Inches(3.8), Inches(4.2))
        c.fill.solid()
        c.fill.fore_color.rgb = PANEL_STRONG if i == 2 else PANEL
        c.line.color.rgb = col_colors[i]
        c.line.transparency = 0.4 if i == 2 else 0.55

        tb(slide, lx + 0.18, 2.14, 3.5, 0.38, yr,  pt=15, color=col_colors[i], bold=True, font="Aptos Display")

        rows = [
            ("Total learners", learners),
            ("Paid users",     paid),
            ("Avg basket",     "₹800"),
            ("ARR",            arr),
        ]
        for j, (label, value) in enumerate(rows):
            row_t = 2.72 + j * 0.8
            tb(slide, lx + 0.18, row_t,        1.5, 0.38, label, pt=11, color=MUTED)
            tb(slide, lx + 1.75, row_t,        1.8, 0.38, value, pt=13, color=TEXT, bold=(label == "ARR"))

    divider(slide, 6.55)
    tb(slide, 0.75, 6.68, 11.85, 0.45,
       "Year 3 ARR of ₹8 Cr is 0.02% of the addressable Indian regional AI education market. Upside is significant.",
       pt=11, color=MUTED, align=PP_ALIGN.CENTER)


def build_team(prs):
    """Slide 9 — Founder / team."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, accent=True)
    title_block(slide, "Founder",
                "Built from lived experience, not theory.",
                sub="Raising seed to move from plan to product.")

    card(slide, 0.75, 2.05, 7.8, 3.6,
         "Rajendar Muddasani — Founder",
         "Telugu-speaking founder with direct experience of the AI knowledge gap in the Indian market. "
         "Building Mitra AI Life Education because no existing platform speaks to the everyday Indian in a "
         "language and format they actually trust.\n\n"
         "Background spans technology, systems thinking, and practical education delivery. "
         "Personal mission: make AI as usable for a homemaker in Warangal as it is for a developer in Bengaluru.",
         strong=True)

    card(slide, 8.75, 2.05, 3.8, 1.7,
         "Why the founder fits",
         [
             "Bilingual — Telugu + English native",
             "Lives in the target market",
             "Content + tech capable",
         ],
         bullet=True, accent_color=CYAN)

    card(slide, 8.75, 3.95, 3.8, 1.7,
         "Planned first hires (post-seed)",
         [
             "Telugu content editor",
             "Visual / comic artist",
             "Part-time community manager",
         ],
         bullet=True, accent_color=GREEN)


def build_ask(prs):
    """Slide 10 — The Ask."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    title_block(slide, "The Ask",
                "₹1 Crore seed round.",
                sub="18-month runway to ₹1 Cr ARR and Series A readiness.")

    card(slide, 0.75, 2.0, 5.75, 3.7,
         "Use of funds",
         [
             "Content & production (L1–L5, Telugu)  —  45%  (₹45L)",
             "Technology & product (site, payments, tracking)  —  28%  (₹28L)",
             "Marketing & community growth  —  18%  (₹18L)",
             "Operations & team  —  9%  (₹9L)",
         ],
         bullet=True, strong=True)

    card(slide, 6.80, 2.0, 5.75, 3.7,
         "Key milestones unlocked by ₹1 Cr",
         [
             "M3 — Levels 1–2 live, 1K beta learners",
             "M6 — Levels 3–5 paid, 10K learners, ₹50L ARR",
             "M12 — 50K learners, ₹1 Cr ARR, Telugu live",
             "M18 — Series A ready: 3x growth, second language",
         ],
         bullet=True, accent_color=GREEN)

    divider(slide, 6.1)
    tb(slide, 0.75, 6.22, 5.75, 0.5,
       "Target raise: ₹1 Cr  |  Instrument: SAFE or equity  |  Stage: Pre-seed / Seed",
       pt=11.5, color=CYAN, bold=True)
    tb(slide, 6.80, 6.22, 5.75, 0.5,
       "mitra-ai-life  |  github.com/rajendarmuddasani/mitra-ai-life",
       pt=11.5, color=MUTED)


# ── main ─────────────────────────────────────────────────────────────────────

def main():
    prs = Presentation()
    prs.slide_width  = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    build_cover(prs)
    build_problem(prs)
    build_solution(prs)
    build_market(prs)
    build_product(prs)
    build_traction(prs)
    build_unit_economics(prs)
    build_revenue(prs)
    build_team(prs)
    build_ask(prs)

    OUTPUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUTPUT_PATH)
    print(f"Wrote {OUTPUT_PATH}")


if __name__ == "__main__":
    main()
